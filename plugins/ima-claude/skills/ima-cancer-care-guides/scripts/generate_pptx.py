"""
Generate a fully branded IMA PowerPoint from a Markdown cancer care guide.

One slide per <!-- pagebreak --> in the source. Cover from YAML front matter.

Usage:
    python3 generate_pptx.py <path.md> [--out <output.pptx>]

Required: pip install python-pptx
"""

import sys, re, yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x00, 0x06, 0x6F)
GOLD      = RGBColor(0xFF, 0xCC, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WARN_BG   = RGBColor(0xFF, 0xF5, 0xCC)
WARN_BDR  = RGBColor(0xFF, 0xCC, 0x00)

# Slide dimensions 16:9
W = Inches(13.33)
H = Inches(7.5)

# ── Low-level helpers ─────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def new_txbox(slide, l, t, w, h, word_wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = word_wrap
    return tb.text_frame


def first_para(tf):
    return tf.paragraphs[0]


def new_para(tf, space_before_pt=0):
    p = tf.add_paragraph()
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    return p


def add_run(para, text, font="Lato", size=12, bold=False, italic=False, color=DARK_GRAY):
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def para_align(para, align):
    para.alignment = align


# ── Inline markdown → runs ────────────────────────────────────────────────────

def apply_inline(para, text, base_size=12, base_color=DARK_GRAY, base_bold=False):
    """Split markdown **bold** and *italic* into separate runs."""
    pattern = re.compile(r'(\*\*.*?\*\*|\*.*?\*)')
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            add_run(para, text[pos:m.start()], size=base_size, bold=base_bold, color=base_color)
        chunk = m.group()
        if chunk.startswith('**'):
            add_run(para, chunk[2:-2], size=base_size, bold=True, color=base_color)
        else:
            add_run(para, chunk[1:-1], size=base_size, italic=True, color=base_color)
        pos = m.end()
    if pos < len(text):
        add_run(para, text[pos:], size=base_size, bold=base_bold, color=base_color)


# ── Markdown parser ───────────────────────────────────────────────────────────

def parse_front_matter(text):
    m = re.match(r'^---\s*\n(.*?)\n---\s*\n', text, re.DOTALL)
    if m:
        return yaml.safe_load(m.group(1)), text[m.end():]
    return {}, text


def parse_blocks(text):
    """Parse markdown into a flat list of typed block dicts."""
    blocks = []
    in_warning = False
    warning_lines = []

    for line in text.splitlines():
        raw = line.rstrip()
        stripped = raw.strip()

        # Warning block
        if stripped == ':::warning':
            in_warning = True
            warning_lines = []
            continue
        if in_warning:
            if stripped == ':::':
                blocks.append({'type': 'warning', 'text': ' '.join(warning_lines)})
                in_warning = False
            else:
                warning_lines.append(stripped)
            continue

        # Page break
        if stripped == '<!-- pagebreak -->':
            blocks.append({'type': 'pagebreak'})
            continue

        # Spacer (skip — not needed in pptx)
        if re.match(r'^<!-- spacer', stripped):
            continue

        # Blockquote / disclaimer
        if stripped.startswith('> '):
            blocks.append({'type': 'disclaimer', 'text': stripped[2:].strip()})
            continue

        # Headings
        if stripped.startswith('#### '):
            blocks.append({'type': 'h4', 'text': stripped[5:].strip()})
            continue
        if stripped.startswith('### '):
            blocks.append({'type': 'h3', 'text': stripped[4:].strip()})
            continue
        if stripped.startswith('## '):
            blocks.append({'type': 'h2', 'text': stripped[3:].strip()})
            continue

        # Numbered list
        m = re.match(r'^(\d+)\.\s+(.*)', stripped)
        if m:
            blocks.append({'type': 'num', 'num': m.group(1), 'text': m.group(2)})
            continue

        # Bullet
        if stripped.startswith('- '):
            blocks.append({'type': 'bullet', 'text': stripped[2:].strip()})
            continue

        # Image (skip in PPTX for now)
        if stripped.startswith('!['):
            blocks.append({'type': 'image_placeholder', 'text': stripped})
            continue

        # Blank line
        if not stripped:
            continue

        # Paragraph
        blocks.append({'type': 'para', 'text': stripped})

    return blocks


def group_into_slides(blocks):
    """Split flat blocks into slide groups at each pagebreak."""
    slides = []
    current = []
    for b in blocks:
        if b['type'] == 'pagebreak':
            if current:
                slides.append(current)
            current = []
        else:
            current.append(b)
    if current:
        slides.append(current)
    return slides


# ── Slide chrome ──────────────────────────────────────────────────────────────

def add_chrome(slide, footer_text=""):
    """Add navy header bar, gold line, footer bar."""
    add_rect(slide, 0, 0, W, Inches(0.5), fill=NAVY)
    add_rect(slide, 0, Inches(0.5), W, Inches(0.03), fill=GOLD)
    add_rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill=NAVY)
    if footer_text:
        tf = new_txbox(slide, Inches(0.4), H - Inches(0.30), Inches(12.5), Inches(0.26))
        p = first_para(tf)
        add_run(p, footer_text, size=8, color=WHITE)


def add_header_text(slide, text):
    tf = new_txbox(slide, Inches(0.4), Inches(0.08), Inches(12), Inches(0.38))
    p = first_para(tf)
    add_run(p, text, size=9, color=WHITE)


# ── Cover slide ───────────────────────────────────────────────────────────────

def build_cover(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=WHITE)

    title1   = meta.get('title1', '').upper()
    title2   = meta.get('title2', '').upper()
    subtitle = meta.get('subtitle', '')
    authors  = meta.get('authors', '')
    date     = meta.get('date', '')

    # Navy hero band
    hero_t = Inches(3.0)
    hero_h = Inches(2.1)
    add_rect(slide, 0, hero_t, W, hero_h, fill=NAVY)

    # Title 1 above hero
    tf = new_txbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.7))
    p = first_para(tf)
    add_run(p, title1, size=80, bold=True, color=DARK_GRAY)

    # Title 2 on hero
    tf2 = new_txbox(slide, Inches(0.8), Inches(3.05), Inches(11.5), Inches(2.0))
    p2 = first_para(tf2)
    add_run(p2, title2, size=72, bold=True, color=WHITE)

    # Subtitle below hero
    tf3 = new_txbox(slide, Inches(0.8), Inches(5.25), Inches(10), Inches(0.55))
    p3 = first_para(tf3)
    add_run(p3, subtitle, size=22, bold=False, color=NAVY)

    # Gold divider
    add_rect(slide, Inches(0.8), Inches(5.9), Inches(8.5), Inches(0.04), fill=GOLD)

    # Authors
    tf4 = new_txbox(slide, Inches(0.8), Inches(6.05), Inches(11), Inches(0.45))
    p4 = first_para(tf4)
    add_run(p4, authors, size=13, bold=True, color=DARK_GRAY)

    # Date
    tf5 = new_txbox(slide, Inches(0.8), Inches(6.55), Inches(8), Inches(0.35))
    p5 = first_para(tf5)
    add_run(p5, date, size=10, color=DARK_GRAY)

    # IMA brand top-right
    tf6 = new_txbox(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.35))
    p6 = first_para(tf6)
    p6.alignment = PP_ALIGN.RIGHT
    add_run(p6, "IMA Health  |  imahealth.org", size=9, color=NAVY)


# ── Content slide ─────────────────────────────────────────────────────────────

CONTENT_L = Inches(0.55)
CONTENT_T = Inches(0.62)
CONTENT_W = Inches(12.2)
CONTENT_H = Inches(6.55)


def render_blocks_to_tf(tf, blocks):
    """Write all blocks into a single text frame."""
    first = True

    def get_para(space=0):
        nonlocal first
        if first:
            first = False
            return first_para(tf)
        return new_para(tf, space_before_pt=space)

    for b in blocks:
        bt = b['type']

        if bt == 'h2':
            p = get_para(12)
            p.alignment = PP_ALIGN.LEFT
            add_run(p, b['text'], size=20, bold=True, color=NAVY)

        elif bt == 'h3':
            p = get_para(8)
            p.alignment = PP_ALIGN.LEFT
            add_run(p, b['text'], size=15, bold=True, color=NAVY)

        elif bt == 'h4':
            p = get_para(6)
            apply_inline(p, b['text'], base_size=13, base_bold=True, base_color=NAVY)

        elif bt == 'para':
            p = get_para(6)
            apply_inline(p, b['text'], base_size=11, base_color=DARK_GRAY)

        elif bt == 'bullet':
            p = get_para(3)
            add_run(p, "•  ", size=11, bold=True, color=NAVY)
            apply_inline(p, b['text'], base_size=11, base_color=DARK_GRAY)

        elif bt == 'num':
            p = get_para(3)
            add_run(p, f"{b['num']}.  ", size=11, bold=True, color=NAVY)
            apply_inline(p, b['text'], base_size=11, base_color=DARK_GRAY)

        elif bt == 'disclaimer':
            p = get_para(6)
            add_run(p, b['text'], size=10, italic=True, color=NAVY)

        elif bt == 'warning':
            p = get_para(8)
            add_run(p, "⚠  " + b['text'], size=11, bold=False, color=RGBColor(0x8B, 0x60, 0x00))

        elif bt == 'image_placeholder':
            p = get_para(6)
            add_run(p, "[Figure]", size=10, italic=True, color=RGBColor(0x99, 0x99, 0x99))


def build_content_slide(prs, blocks, header_text="", footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    add_chrome(slide, footer_text)

    if header_text:
        add_header_text(slide, header_text)

    tf = new_txbox(slide, CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H)
    render_blocks_to_tf(tf, blocks)


# ── Main ──────────────────────────────────────────────────────────────────────

def generate_pptx(md_path, out_path):
    text = Path(md_path).read_text(encoding='utf-8')
    meta, body = parse_front_matter(text)

    blocks     = parse_blocks(body)
    slide_groups = group_into_slides(blocks)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Cover
    build_cover(prs, meta)

    # Footer text for content slides
    title = f"{meta.get('title1','')} {meta.get('title2','')} — {meta.get('date','')}"

    # Content slides — one per pagebreak group
    for group in slide_groups:
        if not group:
            continue
        # Use first h2 in group as header if present
        header = next((b['text'] for b in group if b['type'] == 'h2'), title)
        build_content_slide(prs, group, header_text=header, footer_text=title)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(slide_groups) + 1} slides)")


if __name__ == '__main__':
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument('md')
    ap.add_argument('--out')
    args = ap.parse_args()
    md   = Path(args.md)
    out  = args.out or str(md.with_suffix('.pptx'))
    generate_pptx(str(md), out)
